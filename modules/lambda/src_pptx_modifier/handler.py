"""
EndPoints:
GET  /templates  → lists available .pptx templates from SOURCE_BUCKET
POST /modify     → replaces tokens in a template and saves to OUTPUT_BUCKET
GET  /download   → returns a presigned URL for a processed file
"""

import io
import json
import logging
import os
import re

import boto3
from botocore.exceptions import ClientError
from pptx import Presentation


logger = logging.getLogger()
logger.setLevel(logging.INFO)

SOURCE_BUCKET = os.environ["UNTOUCHED_BUCKET"]
OUTPUT_BUCKET = os.environ["PROCESSED_BUCKET"]
AWS_REGION    = os.environ.get("AWS_REGION", "us-east-1")

s3 = boto3.client("s3", region_name=AWS_REGION)


# ---------------------------------------------------------------------------
# S3
# ---------------------------------------------------------------------------

def list_templates(bucket: str) -> list[dict]:
    response = s3.list_objects_v2(Bucket=bucket)
    return [
        {
            "fileName":    obj["Key"],
            "size":        obj["Size"],
            "lastModified": obj["LastModified"].isoformat(),
        }
        for obj in response.get("Contents", [])
        if obj["Key"].endswith(".pptx")
    ]


def download_pptx(bucket: str, key: str) -> bytes:
    return s3.get_object(Bucket=bucket, Key=key)["Body"].read()


def upload_pptx(bucket: str, key: str, data: bytes) -> None:
    s3.put_object(
        Bucket=bucket,
        Key=key,
        Body=data,
        ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


def build_output_key(original_key: str) -> str:
    return original_key.replace(".pptx", "_modified.pptx")


def generate_presigned_url(bucket: str, key: str, expiration: int = 3600) -> str:
    return s3.generate_presigned_url(
        "get_object",
        Params={"Bucket": bucket, "Key": key},
        ExpiresIn=expiration,
    )


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def build_replacer(replacements: dict):
    pattern = re.compile("|".join(re.escape(k) for k in replacements))
    return lambda text: pattern.sub(lambda m: replacements[m.group(0)], text)


def replace_in_paragraph(paragraph, replacer) -> None:
    if not paragraph.runs:
        return

    full_text = "".join(run.text for run in paragraph.runs).strip()

    if not full_text or "{" not in full_text:
        return

    new_text = replacer(full_text)

    if new_text != full_text:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""


def process_shape(shape, replacer) -> None:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            replace_in_paragraph(paragraph, replacer)

    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    replace_in_paragraph(paragraph, replacer)

    elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub_shape in shape.shapes:
            process_shape(sub_shape, replacer)


def replace_in_presentation(prs: Presentation, replacements: dict) -> None:
    if not replacements:
        return

    replacer = build_replacer(replacements)

    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, replacer)


def modify_pptx(pptx_bytes: bytes, replacements: dict) -> bytes:
    prs = Presentation(io.BytesIO(pptx_bytes))
    replace_in_presentation(prs, replacements)
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


# ---------------------------------------------------------------------------
# Business logic
# ---------------------------------------------------------------------------

def build_replacements(payload: dict) -> dict:
    data = payload["businessData"]
    return {
        "{{PAIN_POINT}}":      data.get("painPoint", ""),
        "{{REVENUE}}":         f"${data.get('revenue', 0):,.0f}",
        "{{ADJUSTED_TARGET}}": f"${round(data.get('revenue', 0) * 0.93):,.0f}",
        "{{TECHNICIANS}}":     str(data.get("technicians", "")),
        "{{REPORTING_DATE}}":  data.get("reportingDate", ""),
    }


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------

def validate_modify(payload: dict) -> str | None:
    if not payload.get("template", {}).get("fileName"):
        return "'template.fileName' is required."
    if not payload["template"]["fileName"].endswith(".pptx"):
        return "'template.fileName' must be a .pptx file."
    if "businessData" not in payload:
        return "'businessData' is required."
    required = ["painPoint", "revenue", "technicians", "reportingDate"]
    missing  = [f for f in required if f not in payload["businessData"]]
    if missing:
        return f"Missing fields in businessData: {', '.join(missing)}."
    return None


# ---------------------------------------------------------------------------
# Route handlers
# ---------------------------------------------------------------------------

def handle_list_templates() -> dict:
    try:
        return success(list_templates(SOURCE_BUCKET))
    except Exception as e:
        logger.exception("Failed to list templates")
        return error(500, str(e))


def handle_modify_pptx(payload: dict) -> dict:
    validation_error = validate_modify(payload)
    if validation_error:
        return error(400, validation_error)

    file_name    = payload["template"]["fileName"]
    replacements = build_replacements(payload)

    try:
        pptx_bytes = download_pptx(SOURCE_BUCKET, file_name)
        modified   = modify_pptx(pptx_bytes, replacements)
        output_key = build_output_key(file_name)

        upload_pptx(OUTPUT_BUCKET, output_key, modified)

        logger.info("Modified '%s' → '%s'", file_name, output_key)
        return success({"outputKey": output_key})

    except ClientError as e:
        if e.response["Error"]["Code"] == "NoSuchKey":
            return error(404, f"Template '{file_name}' not found.")
        logger.exception("S3 error during modify")
        return error(500, str(e))


def handle_download(params: dict) -> dict:
    file_name = params.get("fileName")

    if not file_name:
        return error(400, "Missing 'fileName' query parameter.")

    try:
        s3.head_object(Bucket=OUTPUT_BUCKET, Key=file_name)
        url = generate_presigned_url(OUTPUT_BUCKET, file_name)
        return success({"downloadUrl": url})

    except ClientError as e:
        if e.response["Error"]["Code"] == "404":
            return error(404, f"File '{file_name}' not found.")
        logger.exception("S3 error during download")
        return error(500, str(e))


# ---------------------------------------------------------------------------
# Response
# ---------------------------------------------------------------------------

def success(data) -> dict:
    return {
        "isBase64Encoded": False,
        "statusCode": 200,
        "headers": {
            "Content-Type": "application/json",
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Headers": "Content-Type,X-Api-Key",
            "Access-Control-Allow-Methods": "GET,POST,OPTIONS"
        },
        "body": json.dumps(data),
    }


def error(status_code: int, message: str) -> dict:
    return {
        "isBase64Encoded": False,
        "statusCode": status_code,
        "headers": {
            "Content-Type": "application/json",
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Headers": "Content-Type,X-Api-Key",
            "Access-Control-Allow-Methods": "GET,POST,OPTIONS"
        },
        "body": json.dumps({"error": message}),
    }


# ---------------------------------------------------------------------------
# Router
# ---------------------------------------------------------------------------

def lambda_handler(event, context):
    # Log the incoming event for debugging
    logger.info("Received event: %s", json.dumps(event))

    # Normalize route identification for both API Gateway v1 (REST) and v2 (HTTP)
    http_method = event.get("httpMethod")
    resource = event.get("resource")
    
    if http_method and resource:
        # REST API (v1) format
        route = f"{http_method} {resource}"
    else:
        # HTTP API (v2) format
        route = event.get("routeKey", "")

    # Handle body safely (APIGW v1 can send None for GET requests)
    body = event.get("body")
    if body is None:
        payload = {}
    elif isinstance(body, str):
        try:
            payload = json.loads(body) if body.strip() else {}
        except json.JSONDecodeError:
            payload = {}
    else:
        payload = body

    query_params = event.get("queryStringParameters") or {}

    logger.info("Route detected: %s", route)

    if route == "GET /templates":
        return handle_list_templates()

    if route == "POST /modify":
        return handle_modify_pptx(payload)

    if route == "GET /download":
        return handle_download(query_params)

    return error(404, f"Unknown route: {route}")