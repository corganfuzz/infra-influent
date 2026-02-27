import io
import re

from pptx import Presentation


def build_replacer(replacements: dict):
    pattern = re.compile("|".join(re.escape(k) for k in replacements))
    return lambda text: pattern.sub(lambda m: replacements[m.group(0)], text)


def replace_in_paragraph(paragraph, replacer) -> None:
    if not paragraph.runs:
        return

    full_text = "".join(run.text for run in paragraph.runs).strip()

    if not full_text or "{" not in full_text:
        return

    new_text = replacer(full_text)

    if new_text != full_text:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""


def process_shape(shape, replacer) -> None:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            replace_in_paragraph(paragraph, replacer)

    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    replace_in_paragraph(paragraph, replacer)

    elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub_shape in shape.shapes:
            process_shape(sub_shape, replacer)


def replace_in_presentation(prs: Presentation, replacements: dict) -> None:
    if not replacements:
        return

    replacer = build_replacer(replacements)

    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, replacer)


def modify_pptx(pptx_bytes: bytes, replacements: dict) -> bytes:
    prs = Presentation(io.BytesIO(pptx_bytes))
    replace_in_presentation(prs, replacements)
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def build_replacements(data) -> dict:
    return {
        "{{PAIN_POINT}}":      data.painPoint,
        "{{REVENUE}}":         f"${data.revenue:,.0f}",
        "{{ADJUSTED_TARGET}}": f"${round(data.revenue * 0.93):,.0f}",
        "{{TECHNICIANS}}":     str(data.technicians),
        "{{REPORTING_DATE}}":  data.reportingDate,
    }